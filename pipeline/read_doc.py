"""Extract text from documents (PDF, DOCX, PPTX)."""
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}


def read_document(file_path: str) -> str | None:
    """Extract text from a document file.

    Supports PDF, DOCX, and PPTX. Returns extracted text or None if
    file not found or unsupported format.
    """
    path = Path(file_path)
    if not path.exists():
        return None

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        logger.warning("Unsupported file type: %s", ext)
        return None

    logger.info("Reading %s ...", path.name)

    if ext == ".pdf":
        return _read_pdf(path)
    elif ext == ".docx":
        return _read_docx(path)
    elif ext == ".pptx":
        return _read_pptx(path)


def _read_pdf(path: Path) -> str | None:
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
        result = "\n\n".join(texts).strip()
        logger.info("PDF: extracted %d chars from %d pages", len(result), len(texts))
        return result or None
    except Exception:
        logger.exception("Failed to read PDF: %s", path.name)
        return None


def _read_docx(path: Path) -> str | None:
    try:
        import docx
        doc = docx.Document(str(path))
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        result = "\n\n".join(texts).strip()
        logger.info("DOCX: extracted %d chars", len(result))
        return result or None
    except Exception:
        logger.exception("Failed to read DOCX: %s", path.name)
        return None


def _read_pptx(path: Path) -> str | None:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
        result = "\n\n".join(texts).strip()
        logger.info("PPTX: extracted %d chars from %d slides", len(result), len(prs.slides))
        return result or None
    except Exception:
        logger.exception("Failed to read PPTX: %s", path.name)
        return None
